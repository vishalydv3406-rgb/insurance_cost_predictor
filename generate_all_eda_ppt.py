from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

prs = Presentation()

# Slide 1: Title Slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Insurance AI Pro+"
subtitle.text = "End-to-End Data Science Project\nName: Vishal Yadav\nRegister No: [Replace with your Register No]"

# Slide 2: Abstract and Introduction
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Abstract & Introduction"
content = slide.placeholders[1].text_frame
content.text = "Abstract"
p = content.add_paragraph()
p.text = "Insurance AI Pro+ is an intelligent multi-domain insurance quotation system. It predicts insurance premiums and claim amounts using machine learning while providing Explainable AI (XAI) insights to users."
p.level = 1
p = content.add_paragraph()
p.text = "Introduction"
p.level = 0
p = content.add_paragraph()
p.text = "Traditional actuarial models lack transparency for end-users. This project aims to bridge the gap by deploying dynamic, real-time prediction models across health, life, motor, property, business, specialty, and travel insurance."
p.level = 1

# Slide 3: DataSet detail
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "DataSet Detail"
content = slide.placeholders[1].text_frame
content.text = "Data Characteristics"
p = content.add_paragraph()
p.text = "Source: Custom statistical generation mirroring epidemiological/real-world distributions."
p.level = 1
p = content.add_paragraph()
p.text = "Records: Over 120,000 composite records."
p.level = 1
p = content.add_paragraph()
p.text = "Domains: Claim (50k rows), Health, Life, Motor, Business, Property, Travel, Specialty (each 10k rows)."
p.level = 1
p = content.add_paragraph()
p.text = "Features: Core biological (Age, BMI, Blood Pressure), Categorical (Region, Smoker Status), Financial (Coverage Amount, No_Claim_Years)."
p.level = 1
p = content.add_paragraph()
p.text = "Target Variable: Premium / Claim Amount."
p.level = 1

# Add ALL EDA OUTPUtS dynamically
eda_dir = 'eda_outputs'
if os.path.exists(eda_dir):
    datasets = sorted([d for d in os.listdir(eda_dir) if os.path.isdir(os.path.join(eda_dir, d))])
    for dataset in datasets:
        ds_path = os.path.join(eda_dir, dataset)
        images = sorted([f for f in os.listdir(ds_path) if f.endswith('.png')])
        
        # Group images by pairs (2 per slide)
        for i in range(0, len(images), 2):
            slide_layout = prs.slide_layouts[5] # Title only
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = f"EDA Visualization: {dataset} Dataset (Part {i//2 + 1})"
            
            # Left Image
            img1 = os.path.join(ds_path, images[i])
            slide.shapes.add_picture(img1, Inches(0.5), Inches(1.5), width=Inches(4.3))
            
            # Text label for image 1
            tx1 = slide.shapes.add_textbox(Inches(0.5), Inches(6.0), Inches(4.3), Inches(0.5))
            tx1.text_frame.text = images[i].replace('.png', '').replace('_', ' ')
            tx1.text_frame.paragraphs[0].font.size = Pt(14)
            tx1.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Right Image (if exists)
            if i + 1 < len(images):
                img2 = os.path.join(ds_path, images[i+1])
                slide.shapes.add_picture(img2, Inches(5.2), Inches(1.5), width=Inches(4.3))
                
                # Text label for image 2
                tx2 = slide.shapes.add_textbox(Inches(5.2), Inches(6.0), Inches(4.3), Inches(0.5))
                tx2.text_frame.text = images[i+1].replace('.png', '').replace('_', ' ')
                tx2.text_frame.paragraphs[0].font.size = Pt(14)
                tx2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


# Slide Post EDA: Data Pre Processing & EDA impact
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Data Pre-Processing"
content = slide.placeholders[1].text_frame
content.text = "Transformation Strategies"
p = content.add_paragraph()
p.text = "Missing Values: Median strategy for numericals (SimpleImputer) and Constant strategy ('Unknown') for categoricals."
p.level = 1
p = content.add_paragraph()
p.text = "Outliers: Variables mathematically capped via np.clip to prevent variance skew."
p.level = 1
p = content.add_paragraph()
p.text = "Scaling: Normalization via StandardScaler targeting optimal geometric space for distance algorithms (e.g. KNN/SVR)."
p.level = 1
p = content.add_paragraph()
p.text = "Encoding: LabelEncoder (binary) & OneHotEncoder for multi-class categorical configurations."
p.level = 1
p = content.add_paragraph()
p.text = "Feature Engineering: Synthesized variables via PolynomialFeatures (e.g. Compound actuarial risk of Age + Smoking)."
p.level = 1

# Slide Final: Methodology Planned and Expected Prediction output
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Methodology Built vs Expected"
content = slide.placeholders[1].text_frame
content.text = "Machine Learning Evaluation Baseline:"
p = content.add_paragraph()
p.text = "Linear Algorithms: Linear Regression, SVR, KNN."
p.level = 1
p = content.add_paragraph()
p.text = "Tree Ensembles: Random Forest, Decision Tree."
p.level = 1
p = content.add_paragraph()
p.text = "Boosting Mechanisms: Gradient Boosting, XGBoost."
p.level = 1
p = content.add_paragraph()
p.text = "Final Output Profiles:"
p.level = 0
p = content.add_paragraph()
p.text = "Accuracy Benchmark: Maximum accepted Mean Absolute Error capped safely under ~$550 tolerance."
p.level = 1
p = content.add_paragraph()
p.text = "Best Method Found: Gradient Boosting mapped non-linear incentives successfully (R²: 0.9362 with MAE: $473)."
p.level = 1
p = content.add_paragraph()
p.text = "Risk Profiler: Binary mapping users via Random Forest as High/Low risk yielded a highly resilient 99.71% test accuracy."
p.level = 1

try:
    prs.save('Project_Presentation_Full_EDA.pptx')
    print("Successfully saved Project_Presentation_Full_EDA.pptx")
except Exception as e:
    print(f"Error saving presentation: {e}")
